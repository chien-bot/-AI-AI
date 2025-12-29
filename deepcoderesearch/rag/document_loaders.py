from __future__ import annotations

from pathlib import Path
from typing import Dict, Callable


def _load_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _load_md(path: Path) -> str:
    # Markdown 作为纯文本读取即可，保留标题和段落结构
    return path.read_text(encoding="utf-8", errors="ignore")


def _load_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError as exc:  # pragma: no cover - 依赖缺失提示
        raise RuntimeError("请先安装 pypdf 以支持 PDF 解析: pip install pypdf") from exc

    reader = PdfReader(str(path))
    texts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(texts)


def _load_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("请先安装 python-docx: pip install python-docx") from exc

    doc = docx.Document(str(path))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def _load_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("请先安装 python-pptx: pip install python-pptx") from exc

    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


_LOADERS: Dict[str, Callable[[Path], str]] = {
    ".txt": _load_txt,
    ".md": _load_md,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
}


def load_document(path: Path) -> str:
    """根据扩展名选择合适的 loader。"""

    suffix = path.suffix.lower()
    if suffix not in _LOADERS:
        raise ValueError(f"不支持的文档类型: {suffix}")
    return _LOADERS[suffix](path)
